"""Generate PPTX manuals from markdown slide scripts.

Usage:
  python tools/make_manual_pptx.py

Creates:
  - dist/Glossary_Install_Manual.pptx
  - dist/Glossary_Usage_Manual.pptx

Input:
  - PPT_SLIDES.md
  - PPT_USAGE_MANUAL.md
"""

from __future__ import annotations

from pathlib import Path
import re

from pptx import Presentation
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]


def parse_slides(md: str):
    slides = []
    cur = None
    for line in md.splitlines():
        m = re.match(r"^##\s+슬라이드\s*\d+\.\s*(.+)$", line.strip())
        if m:
            if cur:
                slides.append(cur)
            cur = {"title": m.group(1).strip(), "bullets": [], "notes": []}
            continue
        if cur is None:
            continue
        if line.strip().startswith("-"):
            cur["bullets"].append(line.strip().lstrip("-").strip())
            continue
        if line.strip().startswith("발표자 노트") or line.strip().startswith("[표시]"):
            cur["notes"].append(line.strip())
            continue
        # keep short guidance lines
        if line.strip() and not line.strip().startswith("#"):
            cur["notes"].append(line.strip())

    if cur:
        slides.append(cur)
    return slides


def build_pptx(slides, out_path: Path, title: str):
    prs = Presentation()

    # Title slide
    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    s0.shapes.title.text = title
    s0.placeholders[1].text = "(자동 생성된 매뉴얼 PPT)"

    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = s["title"]
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(s["bullets"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(20)

        # speaker notes
        if s["notes"]:
            notes = slide.notes_slide.notes_text_frame
            notes.text = "\n".join(s["notes"]).strip()[:3000]

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


def main():
    dist = ROOT / "dist"
    dist.mkdir(exist_ok=True)

    install_md = (ROOT / "PPT_SLIDES.md").read_text(encoding="utf-8")
    usage_md = (ROOT / "PPT_USAGE_MANUAL.md").read_text(encoding="utf-8")

    install_slides = parse_slides(install_md)
    usage_slides = parse_slides(usage_md)

    build_pptx(install_slides, dist / "Glossary_Install_Manual.pptx", "제조 AI/DX 용어집 웹앱 - 설치/운영 매뉴얼")
    build_pptx(usage_slides, dist / "Glossary_Usage_Manual.pptx", "제조 AI/DX 용어집 웹앱 - 사용방법 매뉴얼")

    print("Wrote:")
    print(dist / "Glossary_Install_Manual.pptx")
    print(dist / "Glossary_Usage_Manual.pptx")


if __name__ == "__main__":
    main()
